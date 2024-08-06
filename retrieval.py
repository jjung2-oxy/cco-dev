import os
import logging
import traceback
import pandas as pd
import io
import PyPDF2
from pptx import Presentation
from dotenv import load_dotenv
from openai import AzureOpenAI
import numpy as np
from sklearn.metrics.pairwise import cosine_similarity
import json
from azure.storage.blob import BlobServiceClient
from azure.core.exceptions import AzureError
from promptflow.core import tool

# Load environment variables from .env file
load_dotenv()

# Set up logging
log_dir = os.path.join(os.getcwd(), 'logs')
os.makedirs(log_dir, exist_ok=True)
log_file = os.path.join(log_dir, 'retrieval_debug.log')
embeddings_log_file = os.path.join(log_dir, 'embeddings.log')

logging.basicConfig(level=logging.DEBUG, 
                    format='%(asctime)s %(levelname)s: %(message)s',
                    handlers=[logging.FileHandler(log_file),
                              logging.StreamHandler()])

print(f"Debug log file: {log_file}")
print(f"Embeddings log file: {embeddings_log_file}")

# Get Azure Blob Storage connection string from environment variable
connect_str = os.getenv('AZURE_STORAGE_CONNECTION_STRING')
os.environ['AZURE_STORAGE_CONNECTION_STRING'] = connect_str

if not connect_str:
    logging.error("AZURE_STORAGE_CONNECTION_STRING is not set in the environment")
    raise ValueError("AZURE_STORAGE_CONNECTION_STRING is not set")

# Set the container name for blob storage
container_name = "promptflow-storage"

# Set up Azure OpenAI client
client = AzureOpenAI(
    api_key=os.getenv("AZURE_OPENAI_API_KEY"),
    api_version="2023-05-15",
    azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT")
)

def log_embedding(text, embedding):
    with open(embeddings_log_file, "a") as f:
        log_entry = {
            "text": text[:100] + "..." if len(text) > 100 else text,  # Truncate long texts
            "embedding": embedding[:10] + ["..."] if len(embedding) > 10 else embedding,  # Log only first 10 dimensions
            "embedding_length": len(embedding)
        }
        json.dump(log_entry, f)
        f.write("\n")

def get_embedding(text):
    logging.debug(f"Attempting to create embedding with deployment: ada002")
    try:
        response = client.embeddings.create(
            input=text,
            model="ada002"  # Use the actual deployment name from your Azure OpenAI Studio
        )
        logging.debug("Embedding created successfully")
        embedding = response.data[0].embedding
        log_embedding(text, embedding)  # Log the embedding
        return embedding
    except Exception as e:
        logging.error(f"Error creating embedding: {str(e)}")
        raise

@tool
def retrieve_relevant_info(query: str) -> dict:
    logging.debug(f"Received query: {query}")

    if not query:
        logging.info("Empty query, returning default response")
        return {
            "messages": [
                {"role": "system", "content": "You are a helpful AI assistant. Provide a response even if the user query is empty."},
                {"role": "user", "content": "No query provided."},
                {"role": "assistant", "content": "I'm ready to help. What would you like to know about the data in storage?"}
            ],
            "blob_storage_used": False,
            "referenced_files": []
        }

    def read_csv_with_different_encodings(content):
        encodings = ['utf-8', 'iso-8859-1', 'windows-1252', 'utf-16']
        for encoding in encodings:
            try:
                return pd.read_csv(io.BytesIO(content), encoding=encoding)
            except UnicodeDecodeError:
                continue
        raise ValueError(f"Unable to decode CSV with any of the attempted encodings: {', '.join(encodings)}")

    try:
        blob_service_client = BlobServiceClient.from_connection_string(connect_str)
        container_client = blob_service_client.get_container_client(container_name)
        
        all_data = {}
        embeddings = {}
        
        for blob in container_client.list_blobs():
            blob_client = container_client.get_blob_client(blob.name)
            content = blob_client.download_blob().readall()
            
            try:
                if blob.name.lower().endswith('.csv'):
                    df = read_csv_with_different_encodings(content)
                    all_data[blob.name] = f"CSV file with {len(df)} rows and columns: {', '.join(df.columns)}"
                    text_content = df.to_string()
                elif blob.name.lower().endswith('.pdf'):
                    pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
                    all_data[blob.name] = f"PDF file with {len(pdf_reader.pages)} pages"
                    text_content = "\n".join([page.extract_text() for page in pdf_reader.pages])
                elif blob.name.lower().endswith('.pptx'):
                    prs = Presentation(io.BytesIO(content))
                    all_data[blob.name] = f"PowerPoint file with {len(prs.slides)} slides"
                    text_content = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text')])
                elif blob.name.lower().endswith('.txt'):
                    text_content = content.decode('utf-8', errors='ignore')
                    all_data[blob.name] = f"Text file with {len(text_content.split())} words"
                else:
                    all_data[blob.name] = f"Unsupported file type: {blob.name}"
                    continue

                embeddings[blob.name] = get_embedding(text_content[:8000])  # Limit to 8000 characters
            except Exception as e:
                all_data[blob.name] = f"Error processing file: {str(e)}"

        if not all_data:
            logging.error("No files found in the container")
            return {
                "messages": [
                    {"role": "system", "content": "You are a helpful AI assistant. No data files are available."},
                    {"role": "user", "content": query},
                    {"role": "assistant", "content": "I apologize, but there are no data files available in the storage at the moment. Is there anything else I can help you with?"}
                ],
                "blob_storage_used": True,
                "referenced_files": []
            }

        query_embedding = get_embedding(query)
        similarities = {name: cosine_similarity([query_embedding], [emb])[0][0] for name, emb in embeddings.items()}
        most_relevant_file = max(similarities, key=similarities.get)
        
        summary = f"Most relevant file: {most_relevant_file}\n"
        summary += "\n".join([f"{name}: {description}" for name, description in all_data.items()])
        
        logging.debug(f"Data summary: {summary}")

        return {
            "messages": [
                {"role": "system", "content": "You are a helpful AI assistant. Provide concise and accurate responses about the data in storage based on the following summary."},
                {"role": "user", "content": query},
                {"role": "system", "content": f"Data summary: {summary}"},
                {"role": "assistant", "content": "Based on the summary of the data in storage, I'll provide a helpful response."}
            ],
            "blob_storage_used": True,
            "referenced_files": [most_relevant_file]
        }

    except Exception as e:
        logging.error(f"Unexpected error in data retrieval: {str(e)}")
        logging.error(traceback.format_exc())
        return {
            "messages": [
                {"role": "system", "content": "You are a helpful AI assistant. There was an error retrieving information, so provide a general response."},
                {"role": "user", "content": query},
                {"role": "system", "content": f"Error retrieving information: {str(e)}"},
                {"role": "assistant", "content": "I apologize, but I encountered an error while trying to analyze the data in storage. Could you please try your query again or ask about a different aspect of the data?"}
            ],
            "blob_storage_used": True,
            "referenced_files": []
        }

if __name__ == "__main__":
    # Test the get_embedding function
    test_text = "This is a test sentence for embedding."
    try:
        embedding = get_embedding(test_text)
        print(f"Successfully generated embedding of length {len(embedding)}")
        print(f"Embedding has been logged to: {embeddings_log_file}")
    except Exception as e:
        print(f"Error generating embedding: {str(e)}")